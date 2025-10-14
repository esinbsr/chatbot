from pptx import Presentation
from pptx.util import Pt
import re

# export text to pptx
def export_to_pptx(response):
    title_font_size=42
    content_font_size=28
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slides_content = [block.strip() for block in response.split("\n\n") if block.strip()]

    for i, block in enumerate(slides_content, start=1):
        lines = block.split("\n")
        title = lines[0].strip() if lines else f"Diapositive {i}"
        title = re.sub(r'^((Slide\s*\d+\s*.*?(?:-|:)\s*)|(Titre\s*(?:-|:)\s*)|(\d+\s*.))', '', title).strip()
        content_lines = lines[1:] if len(lines) > 1 else []
     
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(title_font_size)

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for line in content_lines:
            line = re.sub(r'^(?:Contenu|Titre)\s*.*?:\s*|^-\s*', '', line).strip()
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(content_font_size)
    prs.save("presentation.pptx")