from pptx import Presentation
from pptx.util import Pt

def export_to_pptx(response):
    title_font_size=36
    content_font_size=20
    prs = Presentation()
    layout = prs.slide_layouts[1]

    slides_content = [block.strip() for block in response.split("\n\n") if block.strip()]

    for i, block in enumerate(slides_content, start=1):
        lines = block.split("\n")
        title = lines[0].strip() if lines else f"Diapositive {i}"
        content_lines = lines[1:] if len(lines) > 1 else []

        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(title_font_size)

        text_frame = slide.placeholders[1].text_frame
        text_frame.clear()

        for line in content_lines:
            if not line.strip():
                continue
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(content_font_size)

    prs.save("presentation.pptx")
